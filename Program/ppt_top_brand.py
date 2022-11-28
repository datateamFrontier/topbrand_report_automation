import pptx
import pandas as pd
import numpy as np
import openpyxl
from pptx import Presentation
from pptx.chart.data import CategoryChartData, XyChartData, ChartData
from pptx.enum.chart import XL_CHART_TYPE
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR, MSO_AUTO_SIZE
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.chart import XL_TICK_LABEL_POSITION, XL_LABEL_POSITION, XL_TICK_MARK,XL_MARKER_STYLE, XL_LEGEND_POSITION
from pptx.oxml.xmlchemy import OxmlElement


def Reverse(lst):
    return [ele for ele in reversed(lst)]

def set_reverse_categories(axis):
    """
    workaround function that replicates the "Categories in Reverse Order" UI option in PPT
    """
    ele = axis._element.xpath(r'c:scaling/c:orientation')[0]
    ele.set("val", "maxMin")
    
def FileTemplate(nama_file_template):
    """
    Mendapatkan file template powerpoint
    
    Parameter
    ---------
    nama_file_template : nama file powerpoint yang akan dipakai sebagai template
    
    Return
    ------
    prs : objek template
    """
    
    prs = Presentation(nama_file_template)
    return prs

def SubElement(parent, tagname, **kwargs):
    """
    Fungsi untuk mengedit kode XML python-pptx pada border tabel
    """
    
    element = OxmlElement(tagname)
    element.attrib.update(kwargs)
    parent.append(element)
    return element

def _set_cell_border(cell, border_color="000000", border_width='12700',
                    lnT=True, lnB=True, lnR=True, lnL=True):
    """ 
    Fungsi untuk mengatur border tabel.
    
    Parameter
    ---------
    border_width : tebal border tabel
    border_color : warna border tabel dengan kode hexadesimal
    lnT : border sebelah atas
    lnB : border sebelah bawah
    lnR : border sebelah kanan
    lnL : border sebelah kiri
    
    Return
    ------
    cell : cell tabel dengan border yang sudah di-edit
    """
    
    tc = cell._tc
    tcPr = tc.get_or_add_tcPr()

    if lnT==True:
        lnT = SubElement(tcPr, 'a:lnT', w=border_width, cap='flat', cmpd='sng', algn='ctr')
        solidFill = SubElement(lnT, 'a:solidFill')
        srgbClr = SubElement(solidFill, 'a:srgbClr', val=border_color)
        prstDash = SubElement(lnT, 'a:prstDash', val='solid')
    
    if lnB==True:
        lnB = SubElement(tcPr, 'a:lnB', w=border_width, cap='flat', cmpd='sng', algn='ctr')
        solidFill = SubElement(lnB, 'a:solidFill')
        srgbClr = SubElement(solidFill, 'a:srgbClr', val=border_color)
        prstDash = SubElement(lnB, 'a:prstDash', val='solid')

    if lnR==True:
        lnR = SubElement(tcPr, 'a:lnR', w=border_width, cap='flat', cmpd='sng', algn='ctr')
        solidFill = SubElement(lnR, 'a:solidFill')
        srgbClr = SubElement(solidFill, 'a:srgbClr', val=border_color)  
        prstDash = SubElement(lnR, 'a:prstDash', val='solid')

    if lnL==True:
        lnL = SubElement(tcPr, 'a:lnL', w=border_width, cap='flat', cmpd='sng', algn='ctr')
        solidFill = SubElement(lnL, 'a:solidFill')
        srgbClr = SubElement(solidFill, 'a:srgbClr', val=border_color)    
        prstDash = SubElement(lnL, 'a:prstDash', val='solid')

    round_ = SubElement(lnB, 'a:round')
    headEnd = SubElement(lnB, 'a:headEnd', type='none', w='med', len='med')
    tailEnd = SubElement(lnB, 'a:tailEnd', type='none', w='med', len='med')
    return cell


def buat_tabel(template_ppt,nomor_slide, 
               baris, kolom, kiri, atas, lebar, tinggi): 
    """
    Membuat objek tabel di power point
    
    Parameter
    ---------
    template_ppt : template ppt yang dipakai
    nomor_slide : angka nomor slide (mulai dari 0)
    baris : jumlah baris tabel
    kolom : jumlah kolom tabel
    kiri : batas kiri tabel dengan tepi slide (dalam satian Inci)
    atas : batas atas tabel dengan tepi slide (dalam satian Inci)
    lebar : ukuran lebar seluruh tabel (dalam satian Inci)
    tinggi : ukuran tinggi seluruh tabel (dalam satian Inci)
    
    Return
    ------
    tabel : objek tabel di ppt
    """
    
    slide = template_ppt.slides[nomor_slide]
    shapes = slide.shapes

    rows = baris
    cols = kolom
    left = Inches(kiri)
    top = Inches(atas)
    width = Inches(lebar)
    height = Inches(tinggi)

    tabel = shapes.add_table(rows, cols, left, top, width, height).table

    return tabel

def plot_text(template_ppt, nomor_slide, teks, nilai, font_style, font_size,
                            kiri, atas, lebar, tinggi):
    """
    Fungsi untuk mem-plot n Sampel dan total multirespon
    
    Parameter
    ---------
    template_ppt : template ppt yang dipakai
    nomor_slide : angka nomor slide (mulai dari 0)
    teks : teks yang akan dimasukkan ke ppt
    nilai : nilai yang akan ditambahkan sebagai teks
    font_style : Style font yang ingin digunakan
    font_size : size font yang akan digunakan
    kiri : batas kiri tabel dengan tepi slide (dalam satian Inci)
    atas : batas atas tabel dengan tepi slide (dalam satian Inci)
    lebar : ukuran lebar seluruh tabel (dalam satian Inci)
    tinggi : ukuran tinggi seluruh tabel (dalam satian Inci)    
    """
    
    left = Inches(kiri)
    top = Inches(atas)
    width = Inches(lebar)
    height = Inches(tinggi)
    
    slide = template_ppt.slides[nomor_slide]
    shapes = slide.shapes
    
    txBox = shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    p = tf.paragraphs[0]
    p.text = teks + ": " + str(round(nilai[0], 0))
 
    tf.paragraphs[0].font.name = font_style
    tf.paragraphs[0].font.size = Pt(font_size)
    tf.paragraphs[0].font.bold = False
    tf.paragraphs[0].alignment = PP_ALIGN.LEFT
    
    return None
    
def plot_nSampel_multirespon(template_ppt, nomor_slide, data_masukan, font_style, font_size,
                            kiri, atas, lebar, tinggi, nS_row, nS_col, 
                            multirespon_ = False, m_row = None, m_col = None):
    """
    Fungsi untuk mem-plot n Sampel dan total multirespon
    
    Parameter
    ---------
    template_ppt : template ppt yang dipakai
    nomor_slide : angka nomor slide (mulai dari 0)
    data_masukan : data yang memuat nilai n Sampel dan multirespon (opsional)
    font_style : Style font yang ingin digunakan
    font_size : size font yang akan digunakan
    kiri : batas kiri tabel dengan tepi slide (dalam satian Inci)
    atas : batas atas tabel dengan tepi slide (dalam satian Inci)
    lebar : ukuran lebar seluruh tabel (dalam satian Inci)
    tinggi : ukuran tinggi seluruh tabel (dalam satian Inci)    
    nS_row : indeks baris di mana nilai n Sampel berada (Mulai dari 0)
    nS_col : indeks kolom di mana nilai n Sampel berada (Mulai dari 0)    
    multirespon_ : True jika ingin menampilkan total multirespon atau False jika tidak
    m_row : indeks baris di mana nilai total multirespon berada (Mulai dari 0)
    m_col : indeks kolom di mana nilai total multirespon berada(Mulai dari 0)
    """
    
    nSampel = data_masukan.iloc[nS_row, nS_col]
    if multirespon_ == True:
        tot_multirespon = data_masukan.iloc[m_row, m_col]
    
    left = Inches(kiri)
    top = Inches(atas)
    width = Inches(lebar)
    height = Inches(tinggi)
    
    slide = template_ppt.slides[nomor_slide]
    shapes = slide.shapes
    
    txBox = shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    p = tf.paragraphs[0]
    p.text = 'n Sampel: ' + str(int(nSampel))
    n_p = 1
    
    if multirespon_ == True:
        p = tf.add_paragraph()
        p.text = 'Total multirespon: ' + str("%.1f" % tot_multirespon) + '%'
        n_p = 2
        
    for i in range(n_p):
        tf.paragraphs[i].font.name = font_style
        tf.paragraphs[i].font.size = Pt(font_size)
        tf.paragraphs[i].font.bold = True
        tf.paragraphs[i].alignment = PP_ALIGN.LEFT
    
    return None
    
def plot_growing_text(template_ppt, nomor_slide, data_plot, list_size, font_style,
                     kiri, atas, lebar, tinggi):
                     
    """
    Fungsi untuk mem-plot teks dengan ukuran naik per paragraf
    
    Parameter
    ---------
    template_ppt : template ppt yang dipakai
    nomor_slide : angka nomor slide (mulai dari 0)
    data_plot : sumber data yang akan diplot isinya
    font_style : Style font yang ingin digunakan
    list_size : size font yang akan digunakan
    kiri : batas kiri tabel dengan tepi slide (dalam satian Inci)
    atas : batas atas tabel dengan tepi slide (dalam satian Inci)
    lebar : ukuran lebar seluruh tabel (dalam satian Inci)
    tinggi : ukuran tinggi seluruh tabel (dalam satian Inci)    
    """
    
    data = data_plot.copy()
    data = data.sort_values(by = 'Presentase', ascending = False)
    teks = data.Kriteria.to_list()
    nilai = data.Presentase.to_list()
    
    left = Inches(kiri)
    top = Inches(atas)
    width = Inches(lebar)
    height = Inches(tinggi)
    
    slide = template_ppt.slides[nomor_slide]
    shapes = slide.shapes
    
    txBox = shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    p = tf.paragraphs[0]
    
    for i in range(len(nilai)):
        if i == 0:
            p.text = str("%.1f" % nilai[i]) + " " + teks[i]
        else:
            p = tf.add_paragraph()
            p.text = str("%.1f" % nilai[i]) + "% " + teks[i]
    
        tf.paragraphs[i].font.name = font_style
        tf.paragraphs[i].font.size = Pt(list_size[i])
        tf.paragraphs[i].font.bold = True
        tf.paragraphs[i].alignment = PP_ALIGN.CENTER
        
def plot_sample_size(data, tabel):
    """
    Fungsi untuk mem-plot tabel sampel size ke ppt
    
    Parameter
    ---------
    data : data yang akan diplot
    tabel : objek tabel yang sudah dibuat
    
    Return
    ------
    None
    """
    header = ['No']+data.columns.to_list()
    for i in range(len(tabel.rows)):
        if i==0:
            list_temp = header
        else:
            if i==len(tabel.rows)-1:
                list_temp = ['']+list(data.loc[i-1,:])
            else:
                list_temp = [i]+list(data.loc[i-1,:])

        for j in range(len(tabel.columns)):     
            if type(list_temp[j])!='str':
                list_temp[j] = str(list_temp[j])

            _set_cell_border(tabel.cell(i,j), "B2B2B2")
            tabel.cell(i,j).text = list_temp[j]
            tabel.cell(i,j).text_frame.paragraphs[0].font.color.rgb = RGBColor(0,0,0)
            tabel.cell(i,j).text_frame.paragraphs[0].font.size = Pt(11)
            tabel.cell(i,j).text_frame.paragraphs[0].font.name = 'Lato'
            tabel.cell(i,j).text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
            tabel.cell(i,j).vertical_anchor = MSO_ANCHOR.MIDDLE
            tabel.cell(i,j).fill.solid()
            tabel.cell(i,j).fill.fore_color.rgb = RGBColor(255, 255, 255) 
            if i == 0 or i == len(tabel.rows) - 1:
                tabel.cell(i,j).text_frame.paragraphs[0].font.color.rgb = RGBColor(255, 255, 255)
                tabel.cell(i,j).text_frame.paragraphs[0].font.bold = True
                tabel.cell(i,j).fill.fore_color.rgb = RGBColor(37, 64, 97)
                if i == len(tabel.rows) - 1 and j != 0 and j != 1:
                    tabel.cell(i,j).text_frame.paragraphs[0].font.bold = False
            if j == 1 and i != 0 and i != len(tabel.rows)-1:
                tabel.cell(i,j).text_frame.paragraphs[0].alignment = PP_ALIGN.LEFT
    tabel.cell(i,0).merge(tabel.cell(i,1))    
    
    return None


def plot_tracking(template_ppt,nomor_slide, 
                  data, merek):
    """
    Plot data tracking
    
    Parameter
    ---------
    template_ppt : template ppt yang dipakai
    nomor_slide : angka nomor slide (mulai dari 0)
    data : data yang akan diploting
    merek : merek tahun ini yang akan diploting
    
    Return
    None
    """
    
    slide = template_ppt.slides[nomor_slide]
    shapes = slide.shapes
    
    gf_tracking = ChartData()
    gf_tracking.categories = data.Tahun.unique()

    for mrk in merek:
        dt_ = data[data.Merek.isin([mrk])].sort_values(by = 'Tahun')
        dt_tbi = dt_['Top Brand Index'].values*100
        tbi = []
        for n_d_tbi in dt_tbi:
            if str(n_d_tbi) == 'nan':
                dt_tbi_i = '#N/A'
            else:
                dt_tbi_i = round(n_d_tbi,1)
            tbi = tbi + [dt_tbi_i]  
        
        gf_tracking.add_series(mrk, tbi)

    x, y, cx, cy = Inches(0.33), Inches(1), Inches(9.26), Inches(4.25)
    chart = slide.shapes.add_chart(
        XL_CHART_TYPE.LINE, x, y, cx, cy, gf_tracking
    ).chart
    chart.has_legend = True
    chart.legend.position = XL_LEGEND_POSITION.BOTTOM
    chart.legend.font.name = 'Tahoma'
    chart.legend.font.size = Pt(12)
    chart.category_axis.has_major_gridlines = False
    chart.value_axis.has_major_gridlines = False

    chart.value_axis.maximum_scale = 30
    chart.value_axis.minimum_scale = 0
    chart.value_axis.visible = False
    
    chart.category_axis.tick_label_position = XL_TICK_LABEL_POSITION.LOW
    chart.category_axis.tick_labels.font.name = 'Tahoma'
    chart.category_axis.tick_labels.font.size = Pt(10)
    
    chart.plots[0].has_data_labels = True
    list_rgb = [[112, 48, 160],[49, 133, 156], [255, 102, 0], [51, 204, 51],[192, 0, 0],[249, 19, 150]]
    for i in range(len(merek)):
        chart.plots[0].series[i].format.line.width = Pt(1)
        chart.plots[0].series[i].format.line.color.rgb = RGBColor(list_rgb[i][0], list_rgb[i][1], list_rgb[i][2])
        dt_ = data[data.Merek == merek[i]].sort_values(by = 'Tahun')
        dt_tbi = dt_['Top Brand Index'].values*100
        dt_tbi = [round(n_d_tbi,1) for n_d_tbi in dt_tbi]
        for j in range(len(dt_tbi)):
            label = dt_tbi[j]
            if str(label) == 'nan':
                chart.plots[0].series[i].points[j].data_label.text_frame.text = ''
            else:
                chart.plots[0].series[i].points[j].data_label.text_frame.text = str(label)+"%"
            #chart.plots[0].series[i].points[j].data_label.ShowBubbleSize = True
            chart.plots[0].series[i].points[j].data_label.font.size = Pt(12)
            chart.plots[0].series[i].points[j].data_label.font.name = 'Tahoma'
            chart.plots[0].series[i].points[j].data_label.font.bold = True
            chart.plots[0].series[i].points[j].data_label.font.color.rgb = RGBColor(list_rgb[i][0], list_rgb[i][1], list_rgb[i][2])
            chart.plots[0].series[i].points[j].data_label.position = XL_LABEL_POSITION.ABOVE
            chart.plots[0].series[i].marker.style = XL_MARKER_STYLE.CIRCLE
            chart.plots[0].series[i].marker.format.fill.solid()
            chart.plots[0].series[i].marker.format.fill.fore_color.rgb = RGBColor(list_rgb[i][0], list_rgb[i][1], list_rgb[i][2])
            chart.plots[0].series[i].marker.size =  8 
    return None


def plot_kemarin_sekarang(template_ppt, nomor_slide,
                          tahun_ini, data_tahun_kemarin, data_tahun_sekarang, merek_top6):
    """
    Plot data tahun kemarin dan sekarang untuk TBI, TOM, LU, FI
    
    Parameter
    ---------
    template_ppt : template ppt yang dipakai
    nomor_slide : angka nomor slide (mulai dari 0)    
    tahun_ini : tahun saat ini
    data_tahun_kemarin : data untuk tahun kemarin
    data_tahun_sekarang : data untuk tahun saat ini
    merek_top6 : daftar merek top6 untuk perbandingan
    
    Return
    ------
    None    
    """
    
    slide = template_ppt.slides[nomor_slide]
    shapes = slide.shapes

    thn = [tahun_ini-1, tahun_ini]
    tampil = ['TOM','LU','FI','TBI']
    list_rgb = [[16,37,63], [23,55,94], [85,142,213], [127, 127, 127]]
    for th_ in range(len(thn)):
        if thn[th_]==tahun_ini:
            data_top6 = data_tahun_sekarang
            x_val = 2.125
        else:
            data_top6 = data_tahun_kemarin
            x_val = 0.125
        for tmpl in range(len(tampil)):
            bar_top6 = CategoryChartData()
            bar_top6.categories = merek_top6
            bar_top6.add_series(tampil[tmpl], 100*data_top6[tampil[tmpl]].values)

            if tampil[tmpl] == 'TOM' and thn[th_] != tahun_ini:
                x, y, cx, cy = Inches(x_val), Inches(1.675), Inches(3.375), Inches(3.5)
            else:
                x, y, cx, cy = Inches(x_val), Inches(1.675), Inches(2.5), Inches(3.5)
                
            bar_top6 = slide.shapes.add_chart(
                XL_CHART_TYPE.BAR_CLUSTERED, x, y, cx, cy, bar_top6
            ).chart

            bar_top6.has_title = False
            bar_top6.has_legend = False
            # tbpd.legend.position = XL_LEGEND_POSITION.BOTTOM
            # tbpd.legend.include_in_layout = False
         
            if tampil[tmpl] != 'TOM' or (tampil[tmpl] == 'TOM' and thn[th_] == tahun_ini):
                bar_top6.category_axis.visible = False
                x_val+=2.125
            else:
                x_val+=3.125
                
            bar_top6.value_axis.minimum_scale = 0
            bar_top6.value_axis.maximum_scale = 100
            bar_top6.value_axis.has_major_gridlines = False
            bar_top6.value_axis.visible = False
            #bar_top6.category_axis.visible = True
            bar_top6.category_axis.tick_labels.font.name = 'Lato'
            bar_top6.category_axis.tick_labels.font.size = Pt(10)
            bar_top6.category_axis.major_tick_mark = XL_TICK_MARK.NONE            
            bar_top6.category_axis.has_minor_gridlines = False
            bar_top6.category_axis.has_major_gridlines = False
            bar_top6.plots[0].gap_width = 50
            set_reverse_categories(bar_top6.category_axis)
            
            bar_top6.plots[0].has_data_labels = True
            data_bar_labels = bar_top6.plots[0].data_labels
            data_bar_labels.number_format = '#.0"%"'
            data_bar_labels.position = XL_LABEL_POSITION.OUTSIDE_END
            data_bar_labels.font.size = Pt(10)
            data_bar_labels.font.name = 'Tahoma'
            
            
            for p in range(6):
                bar_top6.plots[0].series[0].points[p].format.fill.solid()
                bar_top6.plots[0].series[0].points[p].format.fill.fore_color.rgb = RGBColor(list_rgb[tmpl][0], list_rgb[tmpl][1],list_rgb[tmpl][2])
            
    return None


def plot_tracking_client(template_ppt, nomor_slide, data):
    """
    Plot TOM, LU, FI untuk khusus tracking
    
    Parameter
    ---------
    template_ppt : template ppt yang dipakai
    nomor_slide : angka nomor slide (mulai dari 0)
    data : data yang akan di-plot
    
    Return
    ------
    None
    """

    slide = template_ppt.slides[nomor_slide]
    shapes = slide.shapes

    tampil = ['TOM','LU','FI']
    list_th = data.Tahun.to_list()
    list_rgb = [[37, 64, 97],[55, 96, 146], [255, 192, 0]]

    chart = CategoryChartData()
    chart.categories = list_th

    # for thn in range(len(list_th)):
    chart.add_series('FI', 100*data['FI'].values)
    chart.add_series('LU', 100*data['LU'].values)
    chart.add_series('TOM', 100*data['TOM'].values)


    x, y, cx, cy = Inches(0.75), Inches(1.5), Inches(5.44), Inches(5.12)
    chart = slide.shapes.add_chart(
        XL_CHART_TYPE.BAR_CLUSTERED, x, y, cx, cy, chart
    ).chart

    # chart.has_legend = False
    # tbpd.legend.position = XL_LEGEND_POSITION.BOTTOM
    # tbpd.legend.include_in_layout = False

    chart.category_axis.tick_labels.font.size = Pt(12)
    chart.category_axis.tick_labels.font.name = 'Lato Light'
    chart.category_axis.major_tick_mark = XL_TICK_MARK.NONE
    chart.value_axis.major_tick_mark = XL_TICK_MARK.NONE
    chart.value_axis.tick_labels.font.size = Pt(12)
    chart.value_axis.tick_labels.font.name = 'Lato Light' 
    chart.value_axis.tick_labels.number_format = '#.0"%"'
    chart.value_axis.tick_label_position = XL_TICK_LABEL_POSITION.HIGH
    chart.value_axis.has_major_gridlines = True
    chart.value_axis.minimum_scale = 0
    chart.value_axis.maximum_scale = 15
    chart.value_axis.major_unit = 5
    chart.plots[0].gap_width = 125
    chart.plots[0].overlap = -25
        
    chart.plots[0].has_data_labels = True
    data_bar_labels = chart.plots[0].data_labels
    data_bar_labels.number_format = '#.0"%"'
    data_bar_labels.font.name = 'Lato Light'
    data_bar_labels.font.size = Pt(12)
    data_bar_labels.position = XL_LABEL_POSITION.OUTSIDE_END
    
    for i in range(len(tampil)):
        for j in range(len(list_th)):
            chart.plots[0].series[i].points[j].format.fill.solid()
            chart.plots[0].series[i].points[j].format.fill.fore_color.rgb = RGBColor(list_rgb[i][0], list_rgb[i][1],list_rgb[i][2])
 
    return None


def plot_line_tracking_client(template_ppt, nomor_slide, data):
    """
    Plot line chart untuk tracking data client selama 5 tahun terakhir
    
    Parameter
    ---------
    template_ppt : template ppt yang dipakai
    nomor_slide : angka nomor slide (mulai dari 0)
    data : data yang akan di-plot
    
    Return
    ------
    None
    """
    
    kelompok = [['TOM','LU'],['FI','TBI']]
    for i in range(len(kelompok)):
        slide = template_ppt.slides[nomor_slide]
        shapes = slide.shapes

        y_val = 1.25
        for j in range(len(kelompok[i])):
            kel_ = kelompok[i][j]

            gf_client = ChartData()
            gf_client.categories = data.Tahun.to_list()
            gf_client.add_series(kel_, 100*data[kel_].values)

            x, y, cx, cy = Inches(0.75), Inches(y_val), Inches(4.33), Inches(2.5)
            chart = slide.shapes.add_chart(
                XL_CHART_TYPE.LINE_MARKERS, x, y, cx, cy, gf_client
            ).chart
            
            chart.has_title = False
            chart.has_legend = False
    #         chart.legend.position = XL_LEGEND_POSITION.BOTTOM
    
            chart.category_axis.has_major_gridlines = False
            chart.value_axis.has_major_gridlines = False
            chart.value_axis.visible = False
            chart.value_axis.minimum_scale = 0
            chart.value_axis.maximum_scale = 50
            chart.category_axis.tick_labels.font.name = 'Lato'
            chart.category_axis.tick_labels.font.size = Pt(9)
            
            chart.plots[0].series[0].format.line.width = Pt(2.25)
            chart.plots[0].series[0].format.line.color.rgb = RGBColor(0, 32, 96)
            chart.plots[0].series[0].marker.style = XL_MARKER_STYLE.CIRCLE
            chart.plots[0].series[0].marker.format.fill.solid()
            chart.plots[0].series[0].marker.format.fill.fore_color.rgb = RGBColor(0, 32, 96)
            
            chart.plots[0].has_data_labels = True
            data_line_labels = chart.plots[0].data_labels
            data_line_labels.number_format = '#.0"%"'
            data_line_labels.font.name = 'Lato'
            data_line_labels.font.size = Pt(10)
            data_line_labels.position = XL_LABEL_POSITION.ABOVE
            
            y_val+=3
        nomor_slide+=1
        
    return None


def plot_data_ir_kota(template_ppt, nomor_slide, data_ir):
    """
    Plot data ir breakdown kota
    
    Parameter
    ---------
    template_ppt : template ppt yang dipakai
    nomor_slide : angka nomor slide (mulai dari 0)    
    data_ir : data IR yang akan di-plot    
    
    Return
    ------
    None
    """
    
    slide = template_ppt.slides[nomor_slide]
    shapes = slide.shapes

    list_kota = data_ir.kota.to_list()[:len(data_ir)-1]
    list_total = data_ir.Total.to_list()[:len(data_ir)-1]
    list_kategori = []
    for i in range(len(list_total)):
        list_kategori.append(list_kota[i]+', n :'+str(list_total[i]))

    chart_data = CategoryChartData()
    chart_data.categories = list_kategori
    chart_data.add_series('Ya', data_ir.Ya[:len(list_kategori)].values)
    chart_data.add_series('Tidak', data_ir.Tidak[:len(list_kategori)].values)

    x, y, cx, cy = Inches(5), Inches(1.33), Inches(4.75), Inches(4.5)
    chart = slide.shapes.add_chart(
    XL_CHART_TYPE.BAR_STACKED_100, x, y, cx, cy, chart_data
    ).chart
    
    chart.value_axis.visible = False
    chart.value_axis.has_major_gridlines = False
    chart.category_axis.has_major_gridlines = False
    chart.category_axis.tick_labels.font.name = 'Lato Light'
    chart.category_axis.tick_labels.font.size = Pt(10)
    chart.category_axis.major_tick_mark = XL_TICK_MARK.NONE
    set_reverse_categories(chart.category_axis)
    plot = chart.plots[0]
    for i in range(len(list_kota)):
        plot.series[0].points[i].format.fill.solid()
        plot.series[0].points[i].format.fill.fore_color.rgb = RGBColor(69, 114, 167)
        plot.series[1].points[i].format.fill.solid()
        plot.series[1].points[i].format.fill.fore_color.rgb = RGBColor(147, 169, 207)
    
    plot.gap_width = 57
    plot.overlap = 100
    plot.has_data_labels = True
    data_labels = plot.data_labels
    data_labels.number_format = '#.0"%"'
    data_labels.font.name = 'Tahoma (Body)'
    data_labels.font.size = Pt(10)
    data_labels.font.color.rgb = RGBColor(255, 255, 255)
    data_labels.position = XL_LABEL_POSITION.INSIDE_END
    
    return None


def pie_chart_ir(template_ppt, nomor_slide, data_ir, left, top, width, height):
    """
    Membuat chart pie untuk data_ir 'Ya' dan 'Tidak'
    
    Parameter
    ---------
    template_ppt : template ppt yang dipakai
    nomor_slide : angka nomor slide (mulai dari 0)    
    data_ir : data IR yang akan di-plot   
    
    Return
    ------
    None
    """
    slide = template_ppt.slides[nomor_slide]
    shapes = slide.shapes
    
    series = Reverse(data_ir.loc[len(data_ir)-1,:].values[1:3])
    
    pie_chart = ChartData()
    pie_chart.categories = ['Ya','Tidak']
    pie_chart.add_series('IR', series)

    x, y, cx, cy = Inches(left), Inches(top), Inches(width), Inches(height)
    pie= slide.shapes.add_chart(
        XL_CHART_TYPE.PIE_EXPLODED, x, y, cx, cy, pie_chart
    ).chart

    pie.has_title = False
    pie.has_legend = False
    #pie.legend.position = XL_LEGEND_POSITION.BOTTOM
    #pie.legend.include_in_layout = False
    pie.plots[0].gap_width = 0
    pie.plots[0].overlap = 0
    pie.plots[0].has_data_labels = True
    data_pie_labels = pie.plots[0].data_labels
    data_pie_labels.number_format = '#.0"%"'
    data_pie_labels.font.name = 'Tahoma (Body)'
    data_pie_labels.font.size = Pt(20)
    data_pie_labels.font.color.rgb = RGBColor(0,51,102)
    data_pie_labels.position = XL_LABEL_POSITION.OUTSIDE_END
    
    pie.series[0].points[0].format.fill.solid()
    pie.series[0].points[0].format.fill.fore_color.rgb = RGBColor(69, 114, 167)
    pie.series[0].points[1].format.fill.solid()
    pie.series[0].points[1].format.fill.fore_color.rgb = RGBColor(147, 169, 207)
    return None


def donut_chart(template_ppt, nomor_slide, data, 
                     left, top, width, height, explode=False):
    """
    Plot data usia dengan chart dunia
    
    Parameter
    ---------
    template_ppt : template ppt yang dipakai
    nomor_slide : angka nomor slide (mulai dari 0)    
    data : data yang akan di-plot   
    left : tepi chart dari kiri
    top : tepi chart dari atas
    width : lebar chart
    height : tinggi chart    
    explode : jika explode TRUE, chart akan memiliki explosion
    
    Return
    ------
    None
    """

    slide = template_ppt.slides[nomor_slide]
    shapes = slide.shapes

    donut_chart = ChartData()
    donut_chart.categories = Reverse(data.columns.to_list())
    donut_chart.add_series('Series 1', Reverse(data.loc[:0,:].squeeze().values))

    x, y, cx, cy = Inches(left), Inches(top), Inches(width), Inches(height)
    
    if explode==False:
        donut = slide.shapes.add_chart(
            XL_CHART_TYPE.DOUGHNUT, x, y, cx, cy, donut_chart
        ).chart
    else:
        donut = slide.shapes.add_chart(
            XL_CHART_TYPE.DOUGHNUT_EXPLODED, x, y, cx, cy, donut_chart
        ).chart

    donut.has_title = False
    donut.has_legend = False
    #donut.legend.position = XL_LEGEND_POSITION.BOTTOM
    #donut.legend.include_in_layout = False
    donut.plots[0].has_data_labels = True
    data_donut_labels = donut.plots[0].data_labels
    data_donut_labels.number_format = '#.0"%"'
    data_donut_labels.font.name = 'Tahoma (Body)'
    data_donut_labels.font.size = Pt(20)
    data_donut_labels.font.color.rgb = RGBColor(255, 255, 255)
    data_donut_labels.position = XL_LABEL_POSITION.CENTER
    
    list_rgb = [[57, 96, 142], [69, 114, 167], [79, 129, 189], [147, 169, 207], [188, 200, 223]]
    for j in range(len(data.columns)):
        donut.plots[0].series[0].points[j].format.fill.solid()
        donut.plots[0].series[0].points[j].format.fill.fore_color.rgb = RGBColor(list_rgb[j][0], list_rgb[j][1],list_rgb[j][2])
   
    return None


def bar_chart_profil(template_ppt, nomor_slide, data, left, top, width, height, slice_at = None):
    """
    Plot bar chart untuk profil bagian data
    
    Parameter
    ---------
    template_ppt : template ppt yang dipakai
    nomor_slide : angka nomor slide (mulai dari 0)    
    data : data yang akan di-plot   
    left : tepi chart dari kiri
    top : tepi chart dari atas
    width : lebar chart
    height : tinggi chart
    
    Return
    ------
    None
    """
    slide = template_ppt.slides[nomor_slide]
    shapes = slide.shapes
    
    if slice_at != None:
        data = data.sort_values(0, axis = 1, ascending = False)
        data = data.iloc[:,:slice_at]
    
    chart = CategoryChartData()
    chart.categories = data.columns.to_list()
    chart.add_series('Series 1', data.loc[:0,:].squeeze().values)

    x, y, cx, cy = Inches(left), Inches(top), Inches(width), Inches(height)
    chart = slide.shapes.add_chart(
        XL_CHART_TYPE.BAR_CLUSTERED, x, y, cx, cy, chart
    ).chart

    chart.has_title = False
    chart.has_legend = False
    # tbpd.legend.position = XL_LEGEND_POSITION.BOTTOM
    # tbpd.legend.include_in_layout = False

    chart.category_axis.tick_labels.font.size = Pt(10)
    chart.category_axis.tick_labels.font.name = 'Lato Light'
    chart.category_axis.major_tick_mark = XL_TICK_MARK.NONE
    chart.category_axis.has_major_gridlines = False
    chart.value_axis.has_major_gridlines = False
    chart.value_axis.visible = False
    chart.value_axis.minimum_scale = 0
    chart.value_axis.maximum_scale = 70
    chart.plots[0].gap_width = 60
    set_reverse_categories(chart.category_axis)
    
    chart.plots[0].has_data_labels = True
    data_pie_labels = chart.plots[0].data_labels
    data_pie_labels.font.name = 'Tahoma'
    data_pie_labels.font.size = Pt(9)
    data_pie_labels.number_format = '#.0"%"'
    chart.position = XL_LABEL_POSITION.OUTSIDE_END
    
    for j in range(len(data.columns)):
        chart.plots[0].series[0].points[j].format.fill.solid()
        chart.plots[0].series[0].points[j].format.fill.fore_color.rgb = RGBColor(55,96,146)

    return None


def grafik_profil_responden(template_ppt, nomor_slide, data_masukan):
    """
    Membuat bar chart untuk profil responden
    
    Parameter
    ---------
    template_ppt : template ppt yang dipakai
    nomor_slide : angka nomor slide (mulai dari 0)    
    data_masukan : data yang akan di-plot      
    
    Return
    ------
    None
    """
    
    slide = template_ppt.slides[nomor_slide]
    shapes = slide.shapes    
    
    data = data_masukan.copy()
         
    for i in range(len(data)):
        data.iloc[i,1:] = data.iloc[i,1:]/data.iloc[i,-1]*100
    data.iloc[:,1:] = round(data.iloc[:,1:],1)
    
    tampil = data.columns[1:-1].to_list()
    list_th = data[data.columns[0]].to_list()[:-1]
    list_rgb = [[57, 96, 142], [69, 114, 167], [79, 129, 189], [147, 169, 207], [188, 200, 223]]
    chart = CategoryChartData()
    chart.categories = list_th

    for tmp in range(len(tampil)):
        val_ = data[tampil[tmp]].values[:-1]
        chart.add_series(tampil[tmp], val_)

    x, y, cx, cy = Inches(0.33), Inches(1.25), Inches(9.125), Inches(2.5)
    chart = slide.shapes.add_chart(
        XL_CHART_TYPE.COLUMN_CLUSTERED, x, y, cx, cy, chart
    ).chart

    # chart.has_legend = False
    # tbpd.legend.position = XL_LEGEND_POSITION.BOTTOM
    # tbpd.legend.include_in_layout = False

    chart.category_axis.tick_labels.font.size = Pt(10)
    chart.category_axis.tick_labels.font.name = 'Tahoma'
    chart.value_axis.minimum_scale = 0
    chart.value_axis.maximum_scale = 100
    chart.value_axis.visible = False
    chart.value_axis.has_major_gridlines = False
    chart.plots[0].gap_width = 80
    chart.plots[0].overlap = 0
        
    chart.plots[0].has_data_labels = True
    data_bar_labels = chart.plots[0].data_labels
    data_bar_labels.number_format = '#.0"%"'
    data_bar_labels.font.name = 'Tahoma'
    data_bar_labels.font.size = Pt(10)
    data_bar_labels.position = XL_LABEL_POSITION.OUTSIDE_END  
    
    for i in range(len(tampil)):
        for j in range(len(list_th)):
            chart.plots[0].series[i].points[j].format.fill.solid()
            chart.plots[0].series[i].points[j].format.fill.fore_color.rgb = RGBColor(list_rgb[i][0], list_rgb[i][1],list_rgb[i][2])
     
    return None

def plot_tabel_n_sampel(template_ppt, nomor_slide, data_masukan, with_mean = False):
    """
    Membuat plot n sampel
    
    Parameter
    ---------
    template_ppt : template ppt yang dipakai
    nomor_slide : angka nomor slide (mulai dari 0)    
    data_masukan : data yang akan di-plot
    with_mean : True jika ingin menambahkan baris mean (khusus usia)
    
    Return
    ------
    None
    """
    
    data = data_masukan.copy()
    data.set_index('Brand',inplace=True)

    slide = template_ppt.slides[nomor_slide]
    shapes = slide.shapes
    
    if with_mean == False:
        rows = 1
        cols = len(data)-1
        left = Inches(0.5)
        top = Inches(3.75)
        width = Inches(8.875)
        height = Inches(.25)

        tabel_n_sampel = shapes.add_table(rows, cols, left, top, width, height).table
        list_temp = data.iloc[:len(data)-1,-1].to_list()
        
        for j in range(cols):
            list_temp[j] = 'n: ' + str(list_temp[j])
            
            _set_cell_border(tabel_n_sampel.cell(0,j), "969696")
            tabel_n_sampel.cell(0,j).text = list_temp[j]
            tabel_n_sampel.cell(0,j).text_frame.paragraphs[0].font.color.rgb = RGBColor(0,0,0)
            tabel_n_sampel.cell(0,j).text_frame.paragraphs[0].font.name = 'Lato Light'
            tabel_n_sampel.cell(0,j).text_frame.paragraphs[0].font.size = Pt(10)
            tabel_n_sampel.cell(0,j).text_frame.paragraphs[0].font.bold = False
            tabel_n_sampel.cell(0,j).text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
            tabel_n_sampel.cell(0,j).vertical_anchor = MSO_ANCHOR.MIDDLE
            tabel_n_sampel.cell(0,j).margin_left = Inches(0)
            tabel_n_sampel.cell(0,j).margin_right = Inches(0)
            tabel_n_sampel.cell(0,j).fill.solid()
            tabel_n_sampel.cell(0,j).fill.fore_color.rgb = RGBColor(192, 192, 192)
    else:
        usia_ = pd.Series(data.columns[:-1].to_list())
        rata_usia = []
        for i in range(len(data) - 1):
            jumlah = data.iloc[i,:-1].reset_index(drop = True)
            usia_jumlah = round(sum(jumlah*usia_)/data.iloc[i,-1],2)
            rata_usia = rata_usia + [usia_jumlah]
        n_sampel = data.iloc[:-1,-1].to_list()
        
        rows = 2
        cols = len(data)-1
        left = Inches(0.5)
        top = Inches(3.75)
        width = Inches(8.875)
        height = Inches(.625)

        tabel_n_sampel = shapes.add_table(rows, cols, left, top, width, height).table
        list_temp = [rata_usia, n_sampel]
        for i in range(rows):
            for j in range(cols):
                if i == 0:
                    list_temp[i][j] = 'Mean usia:\n' + str("%.2f" % rata_usia[j])
                else:
                    list_temp[i][j] = 'n: ' + str(n_sampel[j])
                    
                _set_cell_border(tabel_n_sampel.cell(i,j), "969696")
                tabel_n_sampel.cell(i,j).text = list_temp[i][j]
                tabel_n_sampel.cell(i,j).vertical_anchor = MSO_ANCHOR.MIDDLE
                tabel_n_sampel.cell(i,j).margin_left = Inches(0)
                tabel_n_sampel.cell(i,j).margin_right = Inches(0)
                tabel_n_sampel.cell(i,j).margin_top = Inches(0)
                tabel_n_sampel.cell(i,j).margin_bottom = Inches(0)
                tabel_n_sampel.cell(i,j).fill.solid()
                tabel_n_sampel.cell(i,j).fill.fore_color.rgb = RGBColor(192, 192, 192)
                tabel_n_sampel.cell(i,j).text_frame.paragraphs[0].font.color.rgb = RGBColor(0,0,0)
                tabel_n_sampel.cell(i,j).text_frame.paragraphs[0].font.name = 'Lato Light'
                tabel_n_sampel.cell(i,j).text_frame.paragraphs[0].font.size = Pt(10)
                tabel_n_sampel.cell(i,j).text_frame.paragraphs[0].font.bold = False
                tabel_n_sampel.cell(i,j).text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
                
                if len(tabel_n_sampel.cell(i,j).text_frame.paragraphs) > 1:
                    tabel_n_sampel.cell(i,j).text_frame.paragraphs[1].font.color.rgb = RGBColor(0,0,0)
                    tabel_n_sampel.cell(i,j).text_frame.paragraphs[1].font.name = 'Lato Light'
                    tabel_n_sampel.cell(i,j).text_frame.paragraphs[1].font.size = Pt(10)
                    tabel_n_sampel.cell(i,j).text_frame.paragraphs[1].font.bold = False
                    tabel_n_sampel.cell(i,j).text_frame.paragraphs[1].alignment = PP_ALIGN.CENTER
                    
        tabel_n_sampel.rows[0].height = Inches(.375)
    return None

def plot_tabel_pekerjaan(template_ppt, nomor_slide, file_excel, nama_sheet):
    """
    Membuat plot tabel pekerjaan
    
    Parameter
    ---------
    template_ppt : template ppt yang dipakai
    nomor_slide : angka nomor slide (mulai dari 0)    
    file_excel = File excel yang memuat data
    nama_sheet = nama sheet data pekerjaan pada file excel
    
    Return
    ------
    None
    """
    
    dt_pekerjaan = pd.read_excel(file_excel, sheet_name=nama_sheet)    
    wb = openpyxl.load_workbook(file_excel)
    sheet_warna = wb[nama_sheet]
       
    slide = template_ppt.slides[nomor_slide]
    shapes = slide.shapes
    rows = len(dt_pekerjaan)+1
    cols = len(dt_pekerjaan.columns)
    left = Inches(0.375)
    top = Inches(1)
    width = Inches(8.125)
    height = Inches(4)

    tabel_pekerjaan = shapes.add_table(rows, cols, left, top, width, height).table

    header = dt_pekerjaan.columns.to_list()

    for i in range(rows):
        if i==0:
            list_temp = header
        else:
            list_temp = dt_pekerjaan.loc[i-1,:].to_list()

        for j in range(cols):     
            if type(list_temp[j])!='str':
                if i != 0 and j != 0 :
                    list_temp[j] = str("%.1f" % list_temp[j]) + '%'
                else:
                    list_temp[j] = str(list_temp[j])
                
            _set_cell_border(tabel_pekerjaan.cell(i,j), "969696")
            tabel_pekerjaan.cell(i,j).text = list_temp[j]
            tabel_pekerjaan.cell(i,j).text_frame.paragraphs[0].font.color.rgb = RGBColor(0,0,0)
            tabel_pekerjaan.cell(i,j).text_frame.paragraphs[0].font.name = 'Lato'
            tabel_pekerjaan.cell(i,j).text_frame.paragraphs[0].font.size = Pt(9)
            tabel_pekerjaan.cell(i,j).text_frame.paragraphs[0].font.bold = False
            tabel_pekerjaan.cell(i,j).text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
            tabel_pekerjaan.cell(i,j).vertical_anchor = MSO_ANCHOR.MIDDLE
            tabel_pekerjaan.cell(i,j).margin_left = Inches(0)
            tabel_pekerjaan.cell(i,j).margin_right = Inches(0)
            tabel_pekerjaan.cell(i,j).fill.solid()
            tabel_pekerjaan.cell(i,j).fill.fore_color.rgb = RGBColor(255, 255, 255)
            
            if i == 0 or i == rows-1:
                tabel_pekerjaan.cell(i,j).text_frame.paragraphs[0].font.color.rgb = RGBColor(255, 255, 255)
                tabel_pekerjaan.cell(i,j).fill.fore_color.rgb = RGBColor(51, 102, 153)  
                
                if j != 0:
                    tabel_pekerjaan.cell(i,j).text_frame.paragraphs[0].font.size = Pt(8)
            else:
                if j != 0:
                    rgb_ij = sheet_warna.cell(i+1, j+1).fill.fgColor.rgb[2:]
                    rgb_ij = tuple(int(rgb_ij[i:i+2], 16) for i in (0, 2, 4))
                    tabel_pekerjaan.cell(i,j).fill.fore_color.rgb = RGBColor(rgb_ij[0], rgb_ij[1], rgb_ij[2]) 
            
            if j == 0: 
                tabel_pekerjaan.cell(i,j).margin_left = Inches(0.1)
                if i != 0:
                    tabel_pekerjaan.cell(i,j).text_frame.paragraphs[0].alignment = PP_ALIGN.LEFT
    #         tabel_pekerjaan.cell(i,j).fill.solid()
    # #         tabel_pekerjaan.cell(i,j).fill.fore_color.rgb = RGBColor(155,187,89)
    # tabel_pekerjaan.cell(i,0).merge(tabel_sampel_size.cell(i,1))            
    tabel_pekerjaan.columns[0].width = Inches(1.75)
    tabel_pekerjaan.rows[0].height = Inches(.5)
    return None

def ppt_plot_bar_gap(template_ppt, nomor_slide, data_plot, gap_data_plot=None, client = 'Comforta'):
    """
    Membuat plot bar chart untuk indeks brand dan gap-nya terhadap indeks tertinggi
    
    Parameter
    ---------
    template_ppt : template ppt yang dipakai
    nomor_slide : angka nomor slide (mulai dari 0)        
    data_plot : data yang akan di-plot (TOM, LU, FI, nonUser, Unaided)
    gap_data_plot : jika ada, tabel data gap indeks antar brand akan di-plot
    client = Nama brand client
    
    Return
    ------
    None
    """
    
    slide = template_ppt.slides[nomor_slide]
    shapes = slide.shapes

    chart = CategoryChartData()
    chart.categories = data_plot.Brand.to_list()
    chart.add_series('Series 1', data_plot[data_plot.columns[1]].values)

    if client == None:
        x, y, cx, cy = Inches(0.25), Inches(.5), Inches(9.25), Inches(2.63)
    else:
        x, y, cx, cy = Inches(0.25), Inches(.5), Inches(9.25), Inches(3)
    chart= slide.shapes.add_chart(
        XL_CHART_TYPE.COLUMN_CLUSTERED, x, y, cx, cy, chart
    ).chart
    
    chart.has_title = False
    chart.has_legend = False
    # tbpd.legend.position = XL_LEGEND_POSITION.BOTTOM
    # tbpd.legend.include_in_layout = False

    chart.category_axis.tick_labels.font.size = Pt(10)
    chart.category_axis.tick_labels.font.name = 'Tahoma'
    chart.category_axis.major_tick_mark = XL_TICK_MARK.NONE
    chart.category_axis.has_major_gridlines = False
    chart.value_axis.has_major_gridlines = False
    chart.value_axis.visible = False
    chart.value_axis.minimum_scale = 0
    chart.value_axis.maximum_scale = 100
    chart.plots[0].gap_width = 50
        
    chart.plots[0].has_data_labels = True
    data_chart_labels = chart.plots[0].data_labels
    data_chart_labels.number_format = '#.0"%"'
    data_chart_labels.font.name = 'Tahoma'
    data_chart_labels.font.size = Pt(10)
    data_chart_labels.position = XL_LABEL_POSITION.OUTSIDE_END  
    
    for j in range(len(data_plot.Brand)):
        chart.plots[0].series[0].points[j].format.fill.solid()
        if client == None:
            chart.category_axis.visible = False
            chart.plots[0].series[0].points[j].format.fill.fore_color.rgb = RGBColor(255, 192, 0)            
        else:
            if data_plot.Brand[j] == client:
                chart.plots[0].series[0].points[j].format.fill.fore_color.rgb = RGBColor(37, 64, 97)
            else:
                chart.plots[0].series[0].points[j].format.fill.fore_color.rgb = RGBColor(127, 127, 127)
    
    
    if type(gap_data_plot)==pd.DataFrame:
        gap_data_plot.columns = ['Brand', 'Gap']
        gap_data_plot = gap_data_plot.iloc[:-1,:]
        rows = 1
        cols = len(gap_data_plot)-2
        left = Inches(1.25)
        top = Inches(4.25)
        width = Inches(7.75)
        height = Inches(.25)

        tabel = shapes.add_table(rows, cols, left, top, width, height).table

        for j in range(cols):
            list_temp = gap_data_plot.loc[j+1,'Gap']

            if type(list_temp)!=str:
                list_temp = str(list_temp) +'%'

            _set_cell_border(tabel.cell(0,j), "969696")
            tabel.cell(0,j).text = list_temp
            tabel.cell(0,j).text_frame.paragraphs[0].font.color.rgb = RGBColor(0,0,0)
            tabel.cell(0,j).text_frame.paragraphs[0].font.size = Pt(10)
            tabel.cell(0,j).text_frame.paragraphs[0].font.name = 'Lato'
            tabel.cell(0,j).text_frame.paragraphs[0].font.bold = False
            tabel.cell(0,j).text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
            tabel.cell(0,j).vertical_anchor = MSO_ANCHOR.MIDDLE
            tabel.cell(0,j).fill.solid()
            tabel.cell(0,j).fill.fore_color.rgb = RGBColor(217, 217, 217)
    return None



def ppt_by_kota(template_ppt, nomor_slide, file_excel, nama_sheet, n_sampel = True, tot_multirespon = False):
    """
    Plot tabel crosstab indeks brand dengan kota
    
    Parameter
    ---------
    template_ppt : template ppt yang dipakai
    nomor_slide : angka nomor slide (mulai dari 0)   
    file_excel = File excel yang memuat data
    nama_sheet = nama sheet data pekerjaan pada file excel    
    
    Return
    ------
    None    
    """
    
    data_kota = pd.read_excel(file_excel, sheet_name=nama_sheet, header=[0,1], index_col=[0])    
    wb = openpyxl.load_workbook(file_excel)
    sheet_warna = wb[nama_sheet]
    slide = template_ppt.slides[nomor_slide]
    shapes = slide.shapes

    rows = data_kota.shape[0]+2
    cols = data_kota.shape[1]+1
    left = Inches(0.125)
    top = Inches(1.125)
    width = Inches(9.25)
    height = Inches(3)

    tabel = shapes.add_table(rows, cols, left, top, width, height).table

    data_kota = data_kota.reset_index()

    isi_header1 = []
    for x in range(len(data_kota.columns)):
        kol = data_kota.columns[x][0]
        if x==0:
            isi_header1.append(kol)
        elif kol==data_kota.columns[x-1][0]:
            isi_header1.append('')
        else:
            isi_header1.append(kol)

    isi_header2 = [data_kota.columns[x][1] for x in range(len(data_kota.columns))]

    for i in range(rows):
        if i==0:
            header = isi_header1
            list_temp = header
        elif i==1:
            header = isi_header2
            list_temp = header

        else:
            list_temp = data_kota.loc[i-2,:].to_list()

        for j in range(cols):     
            if type(list_temp[j])!=str:
                if n_sampel == True:
                    if i != rows - 1:
                        list_temp[j] = str(list_temp[j]) + '%'
                    else:
                        list_temp[j] = int(list_temp[j])
                        list_temp[j] = str(list_temp[j])
                else:
                    list_temp[j] = str(list_temp[j]) + '%'
                
            _set_cell_border(tabel.cell(i,j), "969696")
            tabel.cell(i,j).text = list_temp[j]
            tabel.cell(i,j).text_frame.paragraphs[0].font.color.rgb = RGBColor(0,0,0)
            tabel.cell(i,j).text_frame.paragraphs[0].font.size = Pt(10)
            tabel.cell(i,j).text_frame.paragraphs[0].font.name = 'Lato'
            tabel.cell(i,j).text_frame.paragraphs[0].font.bold = False
            tabel.cell(i,j).text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
            tabel.cell(i,j).vertical_anchor = MSO_ANCHOR.MIDDLE
            tabel.cell(i,j).margin_left = Inches(0)
            tabel.cell(i,j).margin_right = Inches(0)
            tabel.cell(i,j).fill.solid()
            tabel.cell(i,j).fill.fore_color.rgb = RGBColor(255, 255, 255)
            
            if i == 0 or i == 1:
                tabel.cell(i,j).text_frame.paragraphs[0].font.color.rgb = RGBColor(255, 255, 255)
                tabel.cell(i,j).fill.fore_color.rgb = RGBColor(51, 102, 153) 
            else:
                if j != 0:
                    rgb_ij = sheet_warna.cell(i+2, j+1).fill.fgColor.rgb[2:]
                    rgb_ij = tuple(int(rgb_ij[i:i+2], 16) for i in (0, 2, 4))
                    tabel.cell(i,j).fill.fore_color.rgb = RGBColor(rgb_ij[0], rgb_ij[1], rgb_ij[2])                 
            
            if i == 1 and j != 0:
                tabel.cell(i,j).text_frame.paragraphs[0].font.size = Pt(7)
            if j == 0 and i != 0:
                tabel.cell(i,j).text_frame.paragraphs[0].alignment = PP_ALIGN.LEFT
                tabel.cell(i,j).margin_left = Inches(.05)                
            if tot_multirespon == True:
                if i == rows-1 or i == rows-2:
                    tabel.cell(i,j).text_frame.paragraphs[0].font.color.rgb = RGBColor(255, 255, 255)
                    tabel.cell(i,j).fill.fore_color.rgb = RGBColor(51, 102, 153)
                    tabel.cell(i,j).margin_left = Inches(0)
                    tabel.cell(i,j).margin_right = Inches(0)
                if i != 1:
                    tabel.cell(i,j).text_frame.paragraphs[0].font.size = Pt(9)
            elif n_sampel == True and i == rows-1:
                tabel.cell(i,j).text_frame.paragraphs[0].font.color.rgb = RGBColor(255, 255, 255)
                tabel.cell(i,j).fill.fore_color.rgb = RGBColor(51, 102, 153)
                
    tabel.cell(0,0).merge(tabel.cell(1,0))            
    tabel.cell(0,1).merge(tabel.cell(0,7))            
    tabel.cell(0,8).merge(tabel.cell(0,10))     
    tabel.cell(0,11).merge(tabel.cell(0,13))     
    tabel.cell(0,14).merge(tabel.cell(0,15))     
    tabel.columns[0].width = Inches(1)
    
    return None


def ppt_by_usiar(template_ppt, nomor_slide, file_excel, nama_sheet, n_sampel = True, tot_multirespon = False):
    """
    Plot tabel crosstab indeks brand dengan usiar
    
    Parameter
    ---------
    template_ppt : template ppt yang dipakai
    nomor_slide : angka nomor slide (mulai dari 0)  
    file_excel = File excel yang memuat data
    nama_sheet = nama sheet data pekerjaan pada file excel
    
    Return
    ------
    None    
    """
    
    data_usiar = pd.read_excel(file_excel, sheet_name=nama_sheet, index_col = [0])    
    wb = openpyxl.load_workbook(file_excel)
    sheet_warna = wb[nama_sheet]
    
    slide = template_ppt.slides[nomor_slide]
    shapes = slide.shapes

    rows = data_usiar.shape[0]+1
    cols = data_usiar.shape[1]+1
    left = Inches(0.375)
    top = Inches(1.25)
    width = Inches(8.75)
    height = Inches(3)

    tabel = shapes.add_table(rows, cols, left, top, width, height).table

    data_usiar = data_usiar.reset_index()

    isi_header = [data_usiar.columns[x] for x in range(len(data_usiar.columns))]

    for i in range(rows):
        if i==0:
            header = isi_header
            list_temp = header
        else:
            list_temp = data_usiar.loc[i-1,:].to_list()

        for j in range(cols):     
            if type(list_temp[j])!=str:
                if n_sampel == True:
                    if i != rows - 1:
                        list_temp[j] = str(list_temp[j]) + '%'
                    else: 
                        list_temp[j] = int(list_temp[j])
                        list_temp[j] = str(list_temp[j])
                else:
                    list_temp[j] = str(list_temp[j]) + '%'

            _set_cell_border(tabel.cell(i,j), "969696")
            tabel.cell(i,j).text = list_temp[j]
            tabel.cell(i,j).text_frame.paragraphs[0].font.color.rgb = RGBColor(0,0,0)
            tabel.cell(i,j).text_frame.paragraphs[0].font.size = Pt(10)
            tabel.cell(i,j).text_frame.paragraphs[0].font.name = 'Lato'
            tabel.cell(i,j).text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
            #tabel.cell(i,j).text_frame.paragraphs[0].margin_left = Inches(0)
            #tabel.cell(i,j).text_frame.paragraphs[0].margin_right = Inches(0)
            tabel.cell(i,j).vertical_anchor = MSO_ANCHOR.MIDDLE
            tabel.cell(i,j).fill.solid()
            tabel.cell(i,j).fill.fore_color.rgb = RGBColor(255, 255, 255)

            if i == 0:
                tabel.cell(i,j).fill.solid()
                tabel.cell(i,j).fill.fore_color.rgb = RGBColor(51,102,153)                
                tabel.cell(i,j).text_frame.paragraphs[0].font.color.rgb = RGBColor(255,255,255)
            else:
                if j != 0:
                    rgb_ij = sheet_warna.cell(i+1, j+1).fill.fgColor.rgb[2:]
                    rgb_ij = tuple(int(rgb_ij[i:i+2], 16) for i in (0, 2, 4))
                    tabel.cell(i,j).fill.fore_color.rgb = RGBColor(rgb_ij[0], rgb_ij[1], rgb_ij[2])
                    
            if tot_multirespon == True:
                if i == rows-1 or i == rows-2:
                    tabel.cell(i,j).text_frame.paragraphs[0].font.color.rgb = RGBColor(255, 255, 255)
                    tabel.cell(i,j).fill.fore_color.rgb = RGBColor(51, 102, 153)
            elif n_sampel == True and i == rows-1:
                tabel.cell(i,j).text_frame.paragraphs[0].font.color.rgb = RGBColor(255, 255, 255)
                tabel.cell(i,j).fill.fore_color.rgb = RGBColor(51, 102, 153)           
            if j == 0 and i != 0:
                tabel.cell(i,j).text_frame.paragraphs[0].alignment = PP_ALIGN.LEFT
    
    tabel.columns[0].width = Inches(2)
    tabel.rows[0].height = Inches(.5)
    
    return None


def ppt_by_expandr(template_ppt, nomor_slide, file_excel, nama_sheet, n_sampel = True, tot_multirespon = False):
    """
    Plot tabel crosstab indeks brand dengan expandr
    
    Parameter
    ---------
    template_ppt : template ppt yang dipakai
    nomor_slide : angka nomor slide (mulai dari 0) 
    file_excel = File excel yang memuat data
    nama_sheet = nama sheet data pekerjaan pada file excel
    
    Return
    ------
    None    
    """    
 
    data_expandr = pd.read_excel(file_excel, sheet_name=nama_sheet, index_col = [0])    
    wb = openpyxl.load_workbook(file_excel)
    sheet_warna = wb[nama_sheet]
    
    slide = template_ppt.slides[nomor_slide]
    shapes = slide.shapes

    rows = data_expandr.shape[0]+1
    cols = data_expandr.shape[1]+1
    left = Inches(0.5)
    top = Inches(1.125)
    width = Inches(6.25)
    height = Inches(3.25)

    tabel = shapes.add_table(rows, cols, left, top, width, height).table

    data_expandr = data_expandr.reset_index()

    isi_header = [data_expandr.columns[x] for x in range(len(data_expandr.columns))]

    for i in range(rows):
        if i==0:
            header = isi_header
            list_temp = header
        else:
            list_temp = data_expandr.loc[i-1,:].to_list()

        for j in range(cols):     
            if type(list_temp[j])!=str:
                if n_sampel == True:
                    if i != rows - 1:
                        list_temp[j] = str(list_temp[j]) + '%'
                    else: 
                        list_temp[j] = int(list_temp[j])
                        list_temp[j] = str(list_temp[j])
                else:
                    list_temp[j] = str(list_temp[j]) + '%'

            _set_cell_border(tabel.cell(i,j), "969696")
            tabel.cell(i,j).text = list_temp[j]
            tabel.cell(i,j).text_frame.paragraphs[0].font.color.rgb = RGBColor(0,0,0)
            tabel.cell(i,j).text_frame.paragraphs[0].font.size = Pt(10)
            tabel.cell(i,j).text_frame.paragraphs[0].font.name = 'Lato'
            tabel.cell(i,j).text_frame.paragraphs[0].font.bold = False
            tabel.cell(i,j).text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
            tabel.cell(i,j).vertical_anchor = MSO_ANCHOR.MIDDLE
            tabel.cell(i,j).margin_left = Inches(0)
            tabel.cell(i,j).margin_right = Inches(0)
            tabel.cell(i,j).fill.solid()
            tabel.cell(i,j).fill.fore_color.rgb = RGBColor(255, 255, 255)
            
            if i == 0:
                tabel.cell(i,j).fill.solid()
                tabel.cell(i,j).fill.fore_color.rgb = RGBColor(51,102,153)                
                tabel.cell(i,j).text_frame.paragraphs[0].font.color.rgb = RGBColor(255,255,255)
            else:
                if j != 0:
                    rgb_ij = sheet_warna.cell(i+1, j+1).fill.fgColor.rgb[2:]
                    rgb_ij = tuple(int(rgb_ij[i:i+2], 16) for i in (0, 2, 4))
                    tabel.cell(i,j).fill.fore_color.rgb = RGBColor(rgb_ij[0], rgb_ij[1], rgb_ij[2])                
            if tot_multirespon == True:
                if i == rows-1 or i == rows-2:
                    tabel.cell(i,j).text_frame.paragraphs[0].font.color.rgb = RGBColor(255, 255, 255)
                    tabel.cell(i,j).fill.fore_color.rgb = RGBColor(51, 102, 153)
            elif n_sampel == True and i == rows-1:
                tabel.cell(i,j).text_frame.paragraphs[0].font.color.rgb = RGBColor(255, 255, 255)
                tabel.cell(i,j).fill.fore_color.rgb = RGBColor(51, 102, 153)
            if j == 0 and i != 0:
                tabel.cell(i,j).text_frame.paragraphs[0].alignment = PP_ALIGN.LEFT            
    #         tabel.cell(i,j).fill.solid()
    # #         tabel.cell(i,j).fill.fore_color.rgb = RGBColor(155,187,89)
    #         tabel.cell(i,j).text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER

    tabel.columns[0].width = Inches(1.75)
    tabel.rows[0].height = Inches(.5)
    
    return None


def ppt_by_sex(template_ppt, nomor_slide, file_excel, nama_sheet, n_sampel = True, tot_multirespon = False):
    """
    Plot tabel crosstab indeks brand dengan sex
    
    Parameter
    ---------
    template_ppt : template ppt yang dipakai
    nomor_slide : angka nomor slide (mulai dari 0) 
    file_excel = File excel yang memuat data
    nama_sheet = nama sheet data pekerjaan pada file excel
    
    Return
    ------
    None    
    """       

    data_sex = pd.read_excel(file_excel, sheet_name=nama_sheet, index_col = [0])    
    wb = openpyxl.load_workbook(file_excel)
    sheet_warna = wb[nama_sheet]
    
    slide = template_ppt.slides[nomor_slide]
    shapes = slide.shapes
    
    rows = data_sex.shape[0]+1
    cols = data_sex.shape[1]
    left = Inches(7.75)
    top = Inches(1.125)
    width = Inches(1.75)
    height = Inches(3.25)

    tabel1 = shapes.add_table(rows, cols, left, top, width, height).table

    data_sex = data_sex.reset_index()

    isi_header = [data_sex.columns[x] for x in range(len(data_sex.columns))]

    for i in range(rows):
        if i==0:
            header = isi_header
            list_temp = header
        else:
            list_temp = data_sex.loc[i-1,:].to_list()
        list_temp = list_temp[1:]

        for j in range(cols):     
            if type(list_temp[j])!=str:
                if n_sampel == True:
                    if i != rows - 1:
                        list_temp[j] = str(list_temp[j]) + '%'
                    else: 
                        list_temp[j] = int(list_temp[j])
                        list_temp[j] = str(list_temp[j])
                else:
                    list_temp[j] = str(list_temp[j]) + '%'

            _set_cell_border(tabel1.cell(i,j), "969696")
            tabel1.cell(i,j).text = list_temp[j]
            tabel1.cell(i,j).text_frame.paragraphs[0].font.color.rgb = RGBColor(0,0,0)
            tabel1.cell(i,j).text_frame.paragraphs[0].font.size = Pt(10)
            tabel1.cell(i,j).text_frame.paragraphs[0].font.name = 'Lato'
            tabel1.cell(i,j).text_frame.paragraphs[0].font.bold = False
            tabel1.cell(i,j).text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
            tabel1.cell(i,j).vertical_anchor = MSO_ANCHOR.MIDDLE
            tabel1.cell(i,j).fill.solid()
            tabel1.cell(i,j).fill.fore_color.rgb = RGBColor(255, 255, 255)
            
            if i == 0:
                tabel1.cell(i,j).fill.solid()
                tabel1.cell(i,j).fill.fore_color.rgb = RGBColor(51,102,153)                
                tabel1.cell(i,j).text_frame.paragraphs[0].font.color.rgb = RGBColor(255,255,255)
            else:
                rgb_ij = sheet_warna.cell(i+1, j+2).fill.fgColor.rgb[2:]
                rgb_ij = tuple(int(rgb_ij[i:i+2], 16) for i in (0, 2, 4))
                tabel1.cell(i,j).fill.fore_color.rgb = RGBColor(rgb_ij[0], rgb_ij[1], rgb_ij[2])    
                
            if tot_multirespon == True:
                if i == rows-1 or i == rows-2:
                    tabel1.cell(i,j).text_frame.paragraphs[0].font.color.rgb = RGBColor(255, 255, 255)
                    tabel1.cell(i,j).fill.fore_color.rgb = RGBColor(51, 102, 153)
            elif n_sampel == True and i == rows-1:
                tabel1.cell(i,j).text_frame.paragraphs[0].font.color.rgb = RGBColor(255, 255, 255)
                tabel1.cell(i,j).fill.fore_color.rgb = RGBColor(51, 102, 153)    
    #         tabel1.cell(i,j).fill.solid()
    # #         tabel1.cell(i,j).fill.fore_color.rgb = RGBColor(155,187,89)
    #         tabel1.cell(i,j).text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
    
    tabel1.rows[0].height = Inches(.5)
        
    return None


def plot_brand_switching_bar(template_ppt, nomor_slide, data_brandSwitching, tabel_nSample):
    """
    Membuat plot bar chart
    
    Parameter
    ---------
    template_ppt : template ppt yang dipakai
    nomor_slide : angka nomor slide (mulai dari 0)       
    data_brandSwitching : data yang akan diolah
    tabel_nSample : data yang memuat n_Sample
    
    Return
    ------
    None
    """
    
    slide = template_ppt.slides[nomor_slide]
    shapes = slide.shapes

    list_brand = Reverse(data_brandSwitching['Brand'].to_list())
    list_total = Reverse(tabel_nSample['n Sample'].to_list())
    list_total = list_total[1:len(list_total)]
    list_kategori = []
    for i in range(len(list_total)):
        list_kategori.append(list_brand[i]+'\nn :'+str(list_total[i]))

    chart_data = CategoryChartData()
    chart_data.categories = list_kategori
    chart_data.add_series('Loyalist', Reverse(data_brandSwitching['Loyalist'].to_list()))
    chart_data.add_series('Switching out', Reverse(data_brandSwitching['Switching out'].to_list()))

    x, y, cx, cy = Inches(.25), Inches(1.675), Inches(4.25), Inches(3.75)
    chart = slide.shapes.add_chart(
        XL_CHART_TYPE.BAR_STACKED_100, x, y, cx, cy, chart_data
    ).chart

    chart.category_axis.tick_labels.font.name = 'Lato'
    chart.category_axis.tick_labels.font.size = Pt(9)
    chart.value_axis.visible = False
    chart.category_axis.has_major_gridlines = False
    chart.value_axis.has_major_gridlines = False
    chart.category_axis.major_tick_mark = XL_TICK_MARK.NONE
    
    plot = chart.plots[0]
    for i in range(len(list_brand)):
        plot.series[0].points[i].format.fill.solid()
        plot.series[0].points[i].format.fill.fore_color.rgb = RGBColor(69, 114, 167)
        plot.series[1].points[i].format.fill.solid()
        plot.series[1].points[i].format.fill.fore_color.rgb = RGBColor(147, 169, 207)

    plot.gap_width = 54
    plot.overlap = 100
    
    plot.has_data_labels = True    
    data_labels = plot.data_labels
    data_labels.number_format = '#.0"%"'
    data_labels.font.name = 'Tahoma'
    data_labels.font.size = Pt(9)
    data_labels.font.color.rgb = RGBColor(255, 255, 255)
    data_labels.position = XL_LABEL_POSITION.INSIDE_END
    
    return None


def ppt_tabel_brand_switching(template_ppt, nomor_slide, data_brandSwitching,
                             left, top, width, height):
    """
    Menampilkan tabel brand switching analysis
    
    Parameter
    ---------
    template_ppt : template ppt yang dipakai
    nomor_slide : angka nomor slide (mulai dari 0)         
    data_brandSwitching : data yang akan ditampilkan di laporan
    left : tepi chart dari kiri
    top : tepi chart dari atas
    width : lebar chart
    height : tinggi chart    
    
    Return
    ------
    None
    """
    
    slide = template_ppt.slides[nomor_slide]
    shapes = slide.shapes

    rows = data_brandSwitching.shape[0]+2
    cols = data_brandSwitching.shape[1]+1
    left = Inches(left)
    top = Inches(top)
    width = Inches(width)
    height = Inches(height)
    
    tabel = shapes.add_table(rows, cols, left, top, width, height).table
    
    data_brandSwitching = data_brandSwitching.reset_index()
    
    kol_cl = data_brandSwitching.columns.to_list()
    isi_header1 = [kol_cl[i][0] for i in range(len(kol_cl[:2]))]+['']*(len(kol_cl)-2)
    isi_header2 = ['']*1+[kol_cl[i+1][1] for i in range(len(kol_cl[1:]))]

    for i in range(rows):
        if i==0:
            header = isi_header1
            list_temp = header
        elif i==1:
            header = isi_header2
            list_temp = header
        else:
            list_temp = data_brandSwitching.iloc[i-2,:].to_list()
        
        for j in range(cols):     
            #if type(list_temp[j])==float:
                #print(list_temp[j])
            #    list_temp[j]='-'
            if type(list_temp[j])!=str:
                if i == rows-1:
                    list_temp[j] = int(list_temp[j])
                    list_temp[j] = str(list_temp[j])
                else:
                    list_temp[j] = str(list_temp[j]) + '%'
            
            _set_cell_border(tabel.cell(i,j), "969696")
            tabel.cell(i,j).text = list_temp[j]
            tabel.cell(i,j).text_frame.paragraphs[0].font.color.rgb = RGBColor(0,0,0)
            tabel.cell(i,j).text_frame.paragraphs[0].font.size = Pt(10)
            tabel.cell(i,j).text_frame.paragraphs[0].font.name = 'Lato'
            tabel.cell(i,j).text_frame.paragraphs[0].font.bold = False
            tabel.cell(i,j).text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
            tabel.cell(i,j).vertical_anchor = MSO_ANCHOR.MIDDLE
            tabel.cell(i,j).margin_left = Inches(0)
            tabel.cell(i,j).margin_right = Inches(0)
            
            tabel.cell(i,j).fill.solid()
            tabel.cell(i,j).fill.fore_color.rgb = RGBColor(255, 255, 255)
            
            if i == 0 or i== 1 or i == (rows - 1):
                tabel.cell(i,j).text_frame.paragraphs[0].font.color.rgb = RGBColor(255,255,255)
                tabel.cell(i,j).fill.fore_color.rgb = RGBColor(51, 102, 153)
            else:
                if j != 0 and i == j+1:
                    tabel.cell(i,j).text_frame.paragraphs[0].font.color.rgb = RGBColor(255,255,255)
                    tabel.cell(i,j).fill.fore_color.rgb = RGBColor(127, 127, 127)
            if j == 0 and i != 0:
                tabel.cell(i,j).text_frame.paragraphs[0].alignment = PP_ALIGN.LEFT
                tabel.cell(i,j).margin_left = Inches(.05)                       
 
    tabel.cell(0,0).merge(tabel.cell(1,0))            
    tabel.cell(0,1).merge(tabel.cell(0,11))
    
    tabel.columns[0].width = Inches(1.25)
    tabel.rows[0].height = Inches(.25)
    tabel.rows[1].height = Inches(.25)
    return None


def plot_stacked_bar(template_ppt, nomor_slide, list_data, kategori,  kategori1,
                    left, top, width, height):
    """
    Menyajikan stacked_bar chart
    
    Parameter
    ---------
    template_ppt : template ppt yang dipakai
    nomor_slide : angka nomor slide (mulai dari 0)        
    left : tepi chart dari kiri
    top : tepi chart dari atas
    width : lebar chart
    height : tinggi chart        
    list_data : array data series yang akan di tampilkan
    kategori : kategori ini digunakan untuk menamai sumbu dengan nama brand
    kategori1 : kategori ini digunakan untuk menamai sumbu dengan jumlah n_sample
    
    Return
    ------
    None
    """
    
    slide = template_ppt.slides[nomor_slide]
    shapes = slide.shapes

    list_brand = Reverse(kategori.to_list())
    list_total = Reverse(kategori1.to_list())
    list_total = list_total[0:len(list_total)-1]
    list_kategori = []
    for i in range(len(list_total)):
        list_kategori.append(list_brand[i]+'\nn :'+str(list_total[i]))
        
    chart_data = CategoryChartData()
    chart_data.categories = list_kategori
        
    for i in range(len(list_data)):
        chart_data.add_series(list_data[i].name, Reverse(list_data[i].to_list()))
        
    x, y, cx, cy = Inches(left), Inches(top), Inches(width), Inches(height)
    chart = slide.shapes.add_chart(
        XL_CHART_TYPE.BAR_STACKED_100, x, y, cx, cy, chart_data
    ).chart
    
    chart.category_axis.visible = False
    chart.value_axis.visible = False
    chart.category_axis.has_major_gridlines = False
    chart.value_axis.has_major_gridlines = False
    
    plot = chart.plots[0]
    for i in range(len(list_brand)):
        plot.series[0].points[i].format.fill.solid()
        plot.series[0].points[i].format.fill.fore_color.rgb = RGBColor(13, 13, 13)
        plot.series[1].points[i].format.fill.solid()
        plot.series[1].points[i].format.fill.fore_color.rgb = RGBColor(96, 96, 96)
        plot.series[2].points[i].format.fill.solid()
        plot.series[2].points[i].format.fill.fore_color.rgb = RGBColor(146, 208, 80)
    
    plot.gap_width = 54
    plot.overlap = 100
    
    plot.has_data_labels = True
    data_labels = plot.data_labels
    data_labels.number_format = '#.0"%"'
    data_labels.font.name = 'Tahoma'
    data_labels.font.size = Pt(9)
    data_labels.font.color.rgb = RGBColor(255, 255, 255)
    data_labels.position = XL_LABEL_POSITION.CENTER
 
    return None


def plot_bar_chart(template_ppt, nomor_slide, kategori, series, left, top, width,
                   height, min_scale, max_scale, rgb_chart, kategori1=None, label_sumbu = False,
                   font_size = 9, gap_grafik = 54):
    """
    Membuat plot bar chart
    
    Parameter
    ---------
    template_ppt : template ppt yang dipakai
    nomor_slide : angka nomor slide (mulai dari 0)        
    left : tepi chart dari kiri
    top : tepi chart dari atas
    width : lebar chart
    height : tinggi chart      
    kategori : series data untuk kategori di sumbu chart
    series : series data yang akan ditampilkan sebagai chart
    kategori1 : series data untuk kategori tambahan di sumbu chart
    min_scale : minimum scale grafik
    max_scale : maximum scale grafik
    label_sumbu : ya/tidak category axis dimunculkan
    font_size : ukuran font category axis
    gap_grafik : ukuran gap_width grafik
    
    
    Return
    ------
    None    
    """
    
    slide = template_ppt.slides[nomor_slide]
    shapes = slide.shapes
    
    chart = CategoryChartData()
    
    list_kategori =[]
    if type(kategori1)==pd.Series:
        kt_= Reverse(kategori.to_list())
        kt_1 = Reverse(kategori1.to_list())
        kt_1 = kt_1[1:len(kt_1)]
        for i in range(len(kt_1)):
            list_kategori.append(kt_[i]+'\nn :'+str(kt_1[i]))
    else:
        list_kategori = Reverse(kategori.to_list())
        
    chart.categories = list_kategori
    chart.add_series('Series 1', Reverse(series.values))
    
    x, y, cx, cy = Inches(left), Inches(top), Inches(width), Inches(height)
    chart= slide.shapes.add_chart(
        XL_CHART_TYPE.BAR_CLUSTERED, x, y, cx, cy, chart
    ).chart
    
    chart.has_title = False
    chart.has_legend = False
    chart.value_axis.visible = False
    chart.category_axis.tick_labels.font.name = 'Lato'
    chart.category_axis.tick_labels.font.size = Pt(font_size)
    if label_sumbu == False:
        chart.category_axis.visible = False
    chart.value_axis.minimum_scale = min_scale
    chart.value_axis.maximum_scale = max_scale
    chart.value_axis.has_major_gridlines = False
    chart.category_axis.has_major_gridlines = False
    chart.category_axis.major_tick_mark = XL_TICK_MARK.NONE

    # chart.legend.position = XL_LEGEND_POSITION.BOTTOM
    # chart.legend.include_in_layout = False
    
    plot = chart.plots[0]
    for i in range(len(list_kategori)):
        plot.series[0].points[i].format.fill.solid()
        plot.series[0].points[i].format.fill.fore_color.rgb = RGBColor(rgb_chart[0], rgb_chart[1], rgb_chart[2])
    plot.series[0].invert_if_negative = True
    plot.gap_width = gap_grafik
    plot.overlap = 0
    
    plot.has_data_labels = True    
    data_labels = plot.data_labels
    data_labels.number_format = '#.0"%"'
    data_labels.font.name = 'Tahoma'
    data_labels.font.size = Pt(10)
    data_labels.font.color.rgb = RGBColor(0, 0, 0)
    data_labels.position = XL_LABEL_POSITION.OUTSIDE_END
    
    return None


def plot_tabel_brandDiac(template_ppt, nomor_slide, data_tabel):
    """
    Membuat tabel untuk brand Diacnostic
    
    Parameter
    ---------
    template_ppt : template ppt yang dipakai
    nomor_slide : angka nomor slide (mulai dari 0)       
    data_tabel : data yang akan di-plot
    
    Return
    ------
    None
    """

    slide = template_ppt.slides[nomor_slide]
    shapes = slide.shapes

    rows = data_tabel.shape[0]+1
    cols = data_tabel.shape[1]
    left = Inches(.5)
    top = Inches(1)
    width = Inches(6.75)
    height = Inches(3.5)

    tabel = shapes.add_table(rows, cols, left, top, width, height).table

    isi_header = [data_tabel.columns[x] for x in range(len(data_tabel.columns))]

    for i in range(rows):
        if i==0:
            header = isi_header
            list_temp = header
        else:
            list_temp = data_tabel.loc[i-1,:].to_list()

        for j in range(cols):     
            if type(list_temp[j])!=str:
                if j == 1 or j == 2 or j == 3:
                    list_temp[j] = str(list_temp[j]) + '%'
                else:
                    list_temp[j] = str("%.3f" % list_temp[j])
                
            _set_cell_border(tabel.cell(i,j), "969696")
            tabel.cell(i,j).text = list_temp[j]
            tabel.cell(i,j).text_frame.paragraphs[0].font.color.rgb = RGBColor(0,0,0)
            tabel.cell(i,j).text_frame.paragraphs[0].font.size = Pt(10)
            tabel.cell(i,j).text_frame.paragraphs[0].font.name = 'Lato'
            tabel.cell(i,j).text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
            tabel.cell(i,j).vertical_anchor = MSO_ANCHOR.MIDDLE
            
            if i == 0:
                tabel.cell(i,j).text_frame.paragraphs[0].font.color.rgb = RGBColor(255,255,255)
                tabel.cell(i,j).text_frame.paragraphs[0].font.bold = True
                tabel.cell(i,j).fill.solid()
                tabel.cell(i,j).fill.fore_color.rgb = RGBColor(51, 102, 153)
            elif i%2 != 0:
                tabel.cell(i,j).fill.solid()
                tabel.cell(i,j).fill.fore_color.rgb = RGBColor(191, 191, 191)   
            else:
                tabel.cell(i,j).fill.solid()
                tabel.cell(i,j).fill.fore_color.rgb = RGBColor(255, 255, 255)  
            if j == 0 and i != 0:
                tabel.cell(i,j).text_frame.paragraphs[0].alignment = PP_ALIGN.LEFT
    tabel.rows[0].height = Inches(.5)
    tabel.columns[0].width = Inches(1.25)
     
    return None


def plot_scatter_brandDiac(template_ppt, nomor_slide, data_brandDiac, top = None):
    """
    Membuat scatter plot untuk brand diag
    
    Parameter
    ---------
    template_ppt : template ppt yang dipakai
    nomor_slide : angka nomor slide (mulai dari 0)      
    data_brandDiac : data yang akan di-plot
    
    Return
    ------
    None
    """
    slide = template_ppt.slides[nomor_slide]
    shapes = slide.shapes

    chart_data = XyChartData()

    if top != 0:
        data_brandDiac = data_brandDiac.iloc[:top, :]
    
    for i in range(len(data_brandDiac)):
        x = data_brandDiac['TOM/LU'][i]-1
        y = data_brandDiac['FI/LU'][i]-1
        chart_data.add_series(data_brandDiac['Brand'][i]).add_data_point(x,y)
    #'Model %d'%(i)
    x, y, cx, cy = Inches(.75), Inches(2.36), Inches(5.75), Inches(3.75)
    scatter = slide.shapes.add_chart(
        XL_CHART_TYPE.XY_SCATTER, x, y, cx, cy, chart_data
    ).chart

    scatter.has_legend = False
    # scatter.legend.position = XL_LEGEND_POSITION.BOTTOM
    # scatter.legend.include_in_layout = False

    scatter.category_axis.has_major_gridlines = False
    scatter.value_axis.has_major_gridlines = False
    scatter.value_axis.visible = False
    scatter.value_axis.minimum_scale = -1
    scatter.value_axis.maximum_scale = 1
    list_rgb = [[112, 48, 160],[49, 133, 156], [255, 102, 0], [51, 204, 51],[192, 0, 0],[249, 19, 150]]
    
    for i in range(len(scatter.series)):
        scatter.plots[0].series[i].marker.style = XL_MARKER_STYLE.CIRCLE
        scatter.plots[0].series[i].marker.format.fill.solid()
        scatter.plots[0].series[i].marker.format.fill.fore_color.rgb = RGBColor(list_rgb[i][0], list_rgb[i][1], list_rgb[i][2])
        scatter.plots[0].series[i].marker.size = 13    
    #scatter.plots[0].has_data_labels = True 
    #data_label = scatter.plots[0].data_labels
    #data_label.show_series_name = True
        #data_label.text_frame.text = data_brandDiac['Brand'][i]
         
    # data_pie_labels = scatter.plots[0].data_labels
    #data_label.position = XL_LABEL_POSITION.OUTSIDE_END
        
    return None


def ppt_competition_landscape(template_ppt, nomor_slide, file_excel, nama_sheet,
                             left, top, width, height, header_size, font_size,
                             first_row, first_col0, first_col1):
    """
    Membuat tabel competitor landscape di laporan ppt
    
    Parameter
    ---------
    template_ppt : template ppt yang dipakai
    nomor_slide : angka nomor slide (mulai dari 0)        
    left : tepi chart dari kiri
    top : tepi chart dari atas
    width : lebar chart
    height : tinggi chart          
    file_excel = File excel yang memuat data
    nama_sheet = nama sheet data pekerjaan pada file excel
    
    Return
    ------
    None    
    """

    slide = template_ppt.slides[nomor_slide]
    shapes = slide.shapes

    data_cl = pd.read_excel(file_excel, sheet_name=nama_sheet, header = [0,1], index_col = [0,1])    
    wb = openpyxl.load_workbook(file_excel)
    sheet_warna = wb[nama_sheet]
    
    rows = data_cl.shape[0]+2
    cols = data_cl.shape[1]+2
    left = Inches(left)
    top = Inches(top)
    width = Inches(width)
    height = Inches(height)
    
    tabel = shapes.add_table(rows, cols, left, top, width, height).table
    
    data_cl = data_cl.reset_index()
    
    kol_cl = data_cl.columns.to_list()
    isi_header1 = [kol_cl[i][0] for i in range(len(kol_cl[:3]))]+['']*(len(kol_cl)-3)
    isi_header1[1]= ''
    isi_header2 = ['']*2+[kol_cl[i+2][1] for i in range(len(kol_cl[2:]))]

    for i in range(rows):
        if i==0:
            header = isi_header1
            list_temp = header
        elif i==1:
            header = isi_header2
            list_temp = header

        elif i==2:
            list_temp = data_cl.loc[i-2,:].to_list()
        else:
            list_temp = ['']*1+data_cl.iloc[i-2,1:].to_list()

        for j in range(cols):     
            if type(list_temp[j])==float:
                list_temp[j]='-'
            elif type(list_temp[j])!=str:
                list_temp[j] = str("%.3f" % list_temp[j])
            
            _set_cell_border(tabel.cell(i,j), "969696")
            tabel.cell(i,j).text = list_temp[j]
            tabel.cell(i,j).text_frame.paragraphs[0].font.color.rgb = RGBColor(0,0,0)
            tabel.cell(i,j).text_frame.paragraphs[0].font.size = Pt(font_size)
            tabel.cell(i,j).text_frame.paragraphs[0].font.name = 'Lato'
            tabel.cell(i,j).text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
            tabel.cell(i,j).vertical_anchor = MSO_ANCHOR.MIDDLE
            tabel.cell(i,j).margin_left = Inches(0)
            tabel.cell(i,j).margin_right = Inches(0)
            tabel.cell(i,j).margin_top = Inches(0)
            tabel.cell(i,j).margin_bottom = Inches(0)
            tabel.cell(i,j).fill.solid()
            tabel.cell(i,j).fill.fore_color.rgb = RGBColor(255, 255, 255)            
            
            if i == 0 or j == 0:
                tabel.cell(i,j).text_frame.paragraphs[0].font.bold = True
                tabel.cell(i,j).fill.solid()
                tabel.cell(i,j).fill.fore_color.rgb = RGBColor(217, 217, 217)  
            elif i == 1 or j ==1:
                tabel.cell(i,j).text_frame.paragraphs[0].font.color.rgb = RGBColor(255, 255, 255)
                tabel.cell(i,j).fill.solid()
                tabel.cell(i,j).fill.fore_color.rgb = RGBColor(51, 102, 153)                  
                if i == 1:
                    tabel.cell(i,j).text_frame.paragraphs[0].font.size = Pt(header_size)
                else:
                    tabel.cell(i,j).text_frame.paragraphs[0].alignment = PP_ALIGN.LEFT
                    tabel.cell(i,j).margin_left = Inches(0.1)
                    tabel.cell(i,j).margin_right = Inches(0.1)
            else:
                if j != 0 and j != 1:
                    rgb_ij = sheet_warna.cell(i+2, j+1).fill.fgColor.rgb[2:]
                    rgb_ij = tuple(int(rgb_ij[i:i+2], 16) for i in (0, 2, 4))
                    tabel.cell(i,j).fill.fore_color.rgb = RGBColor(rgb_ij[0], rgb_ij[1], rgb_ij[2])   
                    
            if i == j:
                tabel.cell(i,j).text_frame.paragraphs[0].font.color.rgb = RGBColor(51, 102, 153)
                tabel.cell(i,j).fill.solid()
                tabel.cell(i,j).fill.fore_color.rgb = RGBColor(51, 102, 153)

                    
    tabel.cell(0,0).text_frame.paragraphs[0].font.color.rgb = RGBColor(255, 255, 255)
    tabel.cell(0,0).fill.solid()
    tabel.cell(0,0).fill.fore_color.rgb = RGBColor(0, 176, 80)   
    
    tabel.columns[0].width = Inches(first_col0)
    tabel.columns[1].width = Inches(first_col1)
    tabel.rows[0].height = Inches(first_row)
    tabel.rows[1].height = Inches(first_row)
    tabel.cell(0,0).merge(tabel.cell(1,1))            
    tabel.cell(0,2).merge(tabel.cell(0,12))            
    tabel.cell(2,0).merge(tabel.cell(12,0))
     
    return None
    
def pie_chart_olshop(template_ppt, nomor_slide, data_olshop):
    """
    Membuat chart pie untuk data_olshop 'Pernah' dan 'Tidak pernah'
    
    Parameter
    ---------
    template_ppt : template ppt yang dipakai
    nomor_slide : angka nomor slide (mulai dari 0)    
    data_olshop : data olshop yang akan di-plot   
    
    Return
    ------
    None
    """
    slide = template_ppt.slides[nomor_slide]
    shapes = slide.shapes
    
    series = Reverse(data_olshop.loc[len(data_olshop)-1,:].values[1:3])
    
    pie_chart = ChartData()
    pie_chart.categories = ['Pernah','Tidak pernah']
    pie_chart.add_series('olshop', series)

    x, y, cx, cy = Inches(0.5), Inches(2), Inches(4.25), Inches(3.32)
    pie= slide.shapes.add_chart(
        XL_CHART_TYPE.PIE_EXPLODED, x, y, cx, cy, pie_chart
    ).chart

    pie.has_title = False
    pie.has_legend = False
    #pie.legend.position = XL_LEGEND_POSITION.BOTTOM
    #pie.legend.include_in_layout = False
    pie.plots[0].gap_width = 0
    pie.plots[0].overlap = 0
    pie.plots[0].has_data_labels = True
    data_pie_labels = pie.plots[0].data_labels
    data_pie_labels.number_format = '#.0"%"'
    data_pie_labels.font.name = 'Lato light'
    data_pie_labels.font.size = Pt(20)
    data_pie_labels.font.color.rgb = RGBColor(0,51,102)
    data_pie_labels.position = XL_LABEL_POSITION.OUTSIDE_END
    
    pie.series[0].points[0].format.fill.solid()
    pie.series[0].points[0].format.fill.fore_color.rgb = RGBColor(69, 114, 167)
    pie.series[0].points[1].format.fill.solid()
    pie.series[0].points[1].format.fill.fore_color.rgb = RGBColor(147, 169, 207)
    return None
    
def plot_data_olshop_kota(template_ppt, nomor_slide, data_olshop):
    """
    Plot data olshop breakdown kota
    
    Parameter
    ---------
    template_ppt : template ppt yang dipakai
    nomor_slide : angka nomor slide (mulai dari 0)    
    data_olshop : data olshop yang akan di-plot    
    
    Return
    ------
    None
    """
    
    slide = template_ppt.slides[nomor_slide]
    shapes = slide.shapes

    list_kota = data_olshop.kota.to_list()[:len(data_olshop)-1]
    list_total = data_olshop.Total.to_list()[:len(data_olshop)-1]
    list_kategori = []
    for i in range(len(list_total)):
        list_kategori.append(list_kota[i]+', n :'+str(list_total[i]))

    chart_data = CategoryChartData()
    chart_data.categories = list_kategori
    chart_data.add_series('Pernah', data_olshop['Pernah'][:len(list_kategori)].values)
    chart_data.add_series('Tidak pernah', data_olshop['Tidak pernah'][:len(list_kategori)].values)

    x, y, cx, cy = Inches(5), Inches(1), Inches(4.75), Inches(5)
    chart = slide.shapes.add_chart(
    XL_CHART_TYPE.BAR_STACKED_100, x, y, cx, cy, chart_data
    ).chart
    
    chart.value_axis.visible = False
    chart.value_axis.has_major_gridlines = False
    chart.category_axis.has_major_gridlines = False
    chart.category_axis.tick_labels.font.name = 'Lato Light'
    chart.category_axis.tick_labels.font.size = Pt(10)
    chart.category_axis.major_tick_mark = XL_TICK_MARK.NONE
    set_reverse_categories(chart.category_axis)
    plot = chart.plots[0]
    for i in range(len(list_kota)):
        plot.series[0].points[i].format.fill.solid()
        plot.series[0].points[i].format.fill.fore_color.rgb = RGBColor(69, 114, 167)
        plot.series[1].points[i].format.fill.solid()
        plot.series[1].points[i].format.fill.fore_color.rgb = RGBColor(147, 169, 207)
    
    plot.gap_width = 57
    plot.overlap = 100
    plot.has_data_labels = True
    data_labels = plot.data_labels
    data_labels.number_format = '#.0"%"'
    data_labels.font.name = 'Tahoma (Body)'
    data_labels.font.size = Pt(10)
    data_labels.font.color.rgb = RGBColor(255, 255, 255)
    data_labels.position = XL_LABEL_POSITION.INSIDE_END
    
    return None